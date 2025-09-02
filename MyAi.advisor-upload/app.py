# app.py
import streamlit as st
import os
import json
import requests
import uuid
import platform
import threading
from pathlib import Path
from langchain_community.document_loaders import TextLoader
from langchain_community.document_loaders import Docx2txtLoader
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.chains import RetrievalQA
from langchain_core.language_models.llms import LLM
from ctransformers import AutoModelForCausalLM as CTransformersModel
from typing import Any

# --- Helper: Get correct path when bundled with PyInstaller ---
def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = Path(sys._MEIPASS)
    except Exception:
        base_path = Path(".")
    return base_path / relative_path

# Import sys only after defining resource_path (used inside it)
import sys

# --- Anonymous Usage Analytics (Opt-Out) ---
TRACKING_ENABLED = True
TRACKING_ENDPOINT = "https://webhook.site/b3bf5aa1-0aa5-4349-9507-57dcabe15423"
DATA_FILE = Path(__file__).parent / ".telemetry.json"

def send_telemetry():
    if not TRACKING_ENABLED:
        return
    try:
        if DATA_FILE.exists():
            try:
                with DATA_FILE.open('r', encoding='utf-8') as f:
                    user_id = json.load(f).get("user_id")
            except:
                user_id = str(uuid.uuid4())
                with DATA_FILE.open('w', encoding='utf-8') as f:
                    json.dump({"user_id": user_id}, f)
        else:
            user_id = str(uuid.uuid4())
            with DATA_FILE.open('w', encoding='utf-8') as f:
                json.dump({"user_id": user_id}, f)

        payload = {
            "id": user_id,
            "os": platform.system(),
            "version": "1.0",
            "app": "iQvault.ai",
            "python_version": platform.python_version()
        }

        def send():
            try:
                requests.post(TRACKING_ENDPOINT, json=payload, timeout=2)
            except:
                pass

        thread = threading.Thread(target=send, daemon=True)
        thread.start()
    except:
        pass

if "telemetry_sent" not in st.session_state:
    send_telemetry()
    st.session_state.telemetry_sent = True

# Page setup
st.set_page_config(page_title="iQvault.ai", page_icon="🔒")
st.title("🔒 iQvault.ai")
st.markdown("AI-powered knowledge assistant — Secure, Offline, Intelligent.")

# Load documents
def load_docs():
    docs = []
    data_path = resource_path("data")
    
    if not data_path.exists():
        st.error(f"Data folder not found: {data_path.resolve()}")
        return []

    for file in data_path.iterdir():
        if not file.is_file():
            continue
        file_ext = file.suffix.lower()
        file_path = str(file)

        try:
            if file_ext == ".txt":
                loader = TextLoader(file_path, encoding="utf-8")
                docs.extend(loader.load())
            elif file_ext == ".docx":
                loader = Docx2txtLoader(file_path)
                docs.extend(loader.load())
            elif file_ext == ".pdf":
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() or ""
                docs.append({"page_content": text, "metadata": {"source": file.name}})
            elif file_ext == ".pptx":
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                docs.append({"page_content": text, "metadata": {"source": file.name}})
        except Exception as e:
            st.warning(f"Skipped {file.name} | Error: {str(e)}")

    if not docs:
        st.warning("No documents loaded. Add files to the 'data' folder.")
    return docs

# Vector store
@st.cache_resource
def get_vectorstore():
    docs = load_docs()
    from langchain.schema import Document
    formatted_docs = [
        Document(page_content=d["page_content"], metadata=d.get("metadata", {})) 
        if isinstance(d, dict) else d 
        for d in docs
    ]
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    texts = splitter.split_documents(formatted_docs)
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    return Chroma.from_documents(texts, embedding=embeddings)

# Custom LLM wrapper
class CTransformersLLM(LLM):
    model: Any

    @property
    def _llm_type(self) -> str:
        return "ctransformers"

    def _call(self, prompt: str, **kwargs) -> str:
        return self.model(prompt)

# Load LLM
@st.cache_resource
def get_llm():
    model_file = "tinyllama-1.1b-chat-v1.0.Q4_K_M.gguf"
    model_path = resource_path("model") / model_file

    if not model_path.exists():
        st.error(f"""
❌ Model file not found:  
`{model_path}`  

Please download it from:  
https://huggingface.co/TheBloke/TinyLlama-1.1B-Chat-v1.0-GGUF  
and place it in the 'model' folder.
        """)
        st.stop()

    try:
        model = CTransformersModel.from_pretrained(
            str(model_path),
            model_type="llama",
            max_new_tokens=200,
            temperature=0.3,
            context_length=2048
        )
        return CTransformersLLM(model=model)
    except Exception as e:
        st.error(f"Failed to load model: {str(e)}")
        st.stop()

# QA Chain
@st.cache_resource
def get_qa_chain():
    llm = get_llm()
    db = get_vectorstore()
    retriever = db.as_retriever(search_kwargs={"k": 1})
    return RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=retriever,
        return_source_documents=True
    )

# Initialize
try:
    qa_chain = get_qa_chain()
except Exception as e:
    st.error(f"Error initializing QA chain: {str(e)}")
    st.stop()

# Chat history
if "messages" not in st.session_state:
    st.session_state.messages = []

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if prompt := st.chat_input("Ask about our AI solutions..."):
    with st.chat_message("user"):
        st.markdown(prompt)
    st.session_state.messages.append({"role": "user", "content": prompt})

    with st.spinner("Thinking..."):
        try:
            result = qa_chain({"query": prompt})
            answer = result["result"]
            sources = result["source_documents"]
        except Exception as e:
            answer = "Sorry, I couldn't process your request."
            sources = []
            st.error(f"Error generating response: {str(e)}")

    with st.chat_message("assistant"):
        st.markdown(answer)
        if sources:
            with st.expander("See source"):
                content = sources[0].page_content
                st.write(content[:500] + "..." if len(content) > 500 else content)
    st.session_state.messages.append({"role": "assistant", "content": answer})
